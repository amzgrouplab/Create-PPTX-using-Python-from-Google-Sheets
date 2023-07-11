import sys
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE_TYPE
if __name__ == '__main__':
    file_name = sys.argv[1]

    print("the file " + file_name)

    slide_title = pd.read_csv(file_name)
    slide_title = slide_title.columns[0]
    df = pd.read_csv(file_name, header=1)

    df = df[~df.iloc[:, 0].str.contains('^\\*BlankRow')]

    # Reset the index of the DataFrame
    df.reset_index(drop=True, inplace=True)

    existing_pptx = "./slide-template-empty.pptx"
    prs = Presentation(existing_pptx)

    # Get the slide layout from the existing presentation
    slide_layout = prs.slides[0].slide_layout

    slide = prs.slides.add_slide(slide_layout)
  
    
    # Read slide-specifications.csv file
    specifications_df = pd.read_csv('slide-specifications.csv')
    print (specifications_df)

    # Extract relevant information
    title_start = (specifications_df.columns[0], specifications_df.columns[1])
    table_start = (specifications_df.columns[1], specifications_df.columns[1])

    title_x = float(specifications_df.columns[0])
    title_y = float(specifications_df.columns[1])
    # Add a title shape to the slide
    title_placeholder = slide.shapes.title
    if not title_placeholder:
        title_placeholder = slide.shapes.add_shape(
            autoshape_type_id=1,
            left=Inches(title_x),
            top=Inches(title_y),
            width=Inches(8),
            height=Inches(1)
        )
    title_placeholder.text = slide_title

    # Set the font size for the title
    font = title_placeholder.text_frame.paragraphs[0].runs[0].font
    font.size = Pt(20)  # Specify the desired font size in points (e.g., 20)

    # Move the title to the top left corner
    title_placeholder.left = Inches(0.5)
    title_placeholder.top = Inches(0.5)

    # Adjust the width of the title shape
    title_placeholder.width = Inches(2)  # Specify the desired width in inches (e.g., 5)

    # Set the text direction to horizontal
    text_frame = title_placeholder.text_frame
    text_frame.text = title_placeholder.text  # Refresh the text to apply the changes

    # Set the rotation transformation for horizontal text
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = run.text.strip()  # Remove leading/trailing spaces to avoid potential line breaks

    text_frame.rotation = 0

    # Set the alignment of the text frame to left
    text_frame.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # Adjust autofit properties
    text_frame.auto_size = True

    # Set the font size for the title (alternative approach)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)  # Specify the desired font size in points (e.g., 20)

    # Add a table shape to the slide
    left = Inches(0)
    top = Inches(2)
    width = Inches(10)
    height = Inches(3)
    table = slide.shapes.add_table(rows=df.shape[0] + 1, cols=df.shape[1], left=left, top=top, width=width,
                                   height=height).table

    # Remove bullet points
    for cell in table.iter_cells():
        for paragraph in cell.text_frame.paragraphs:
            paragraph.clear()

    # Set the column names and font size
    for i, column_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column_name
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(15)

    # Set the values from the dataframe and font size
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(15)

    # Auto-fit column widths
    for column in table.columns:
        column.width = int(width / df.shape[1])

    prs.save("presentation.pptx")